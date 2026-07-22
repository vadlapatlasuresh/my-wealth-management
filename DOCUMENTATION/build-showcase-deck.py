#!/usr/bin/env python3
"""TerraVest showcase deck — complete build.

Contents:
  Part 1  App overview + full competitive analysis (pricing, 2 capability matrices,
          stack-cost, our precise offering)
  Part 2  All 46 features, one slide each, demo steps + say-this panel
  Part 3  Presenter notes (transitions, pacing, moments, closes, Q&A)
  Part 4  Running the demo well + open items + cheat card

Competitor pricing verified July 2026 via web research.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------------- brand
FOREST      = RGBColor(0x1A, 0x4D, 0x3B)
FOREST_DEEP = RGBColor(0x10, 0x33, 0x27)
FOREST_MID  = RGBColor(0x2D, 0x6B, 0x52)
GOLD        = RGBColor(0xC9, 0x97, 0x3A)
GOLD_LIGHT  = RGBColor(0xF0, 0xC8, 0x78)
CREAM       = RGBColor(0xFA, 0xF8, 0xF3)
INK         = RGBColor(0x1C, 0x1C, 0x1A)
SLATE       = RGBColor(0x5A, 0x5F, 0x5B)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
SAGE        = RGBColor(0xE4, 0xEC, 0xE6)
MIST        = RGBColor(0xE8, 0xF0, 0xE9)
GREY        = RGBColor(0xB8, 0xBD, 0xB9)

W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def bg(slide, color):
    r = slide.shapes.add_shape(1, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    return r


def band(slide, x, y, w, h, color):
    r = slide.shapes.add_shape(1, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    return r


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line=1.25, space_after=0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(runs):
        s, size, bold, color = item[:4]
        sb = item[4] if len(item) > 4 else None
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line
        if space_after:
            p.space_after = Pt(space_after)
        if sb:
            p.space_before = Pt(sb)
        r = p.add_run(); r.text = s
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = "Helvetica Neue"
    return tb


def notes(slide, body):
    slide.notes_slide.notes_text_frame.text = body


def rule(slide, x, y, w, color=GOLD, h=Inches(0.045)):
    return band(slide, x, y, w, h, color)


# ---------------------------------------------------------------- slide types
def slide_title(title, subtitle, kicker=None, footer=None, note=""):
    s = prs.slides.add_slide(BLANK)
    bg(s, FOREST_DEEP)
    band(s, 0, 0, Inches(0.28), H, GOLD)
    band(s, Inches(9.6), Inches(0), Inches(3.733), H, FOREST)
    if kicker:
        text(s, Inches(1.15), Inches(1.5), Inches(8.2), Inches(0.4),
             [(kicker.upper(), 13, True, GOLD_LIGHT)])
    text(s, Inches(1.15), Inches(2.05), Inches(8.3), Inches(1.8),
         [(title, 50, True, WHITE)], line=1.06)
    rule(s, Inches(1.15), Inches(4.02), Inches(1.5))
    text(s, Inches(1.15), Inches(4.35), Inches(8.0), Inches(1.4),
         [(subtitle, 19, False, SAGE)], line=1.35)
    if footer:
        text(s, Inches(1.15), Inches(6.5), Inches(9.0), Inches(0.5),
             [(footer, 12, False, RGBColor(0x8F, 0xA5, 0x99))])
    if note:
        notes(s, note)
    return s


def slide_section(num, title, subtitle, note=""):
    s = prs.slides.add_slide(BLANK)
    bg(s, FOREST)
    band(s, 0, 0, Inches(0.28), H, GOLD)
    text(s, Inches(1.15), Inches(2.35), Inches(2.0), Inches(1.2),
         [(num, 82, True, RGBColor(0x3E, 0x7E, 0x64))], line=1.0)
    text(s, Inches(1.15), Inches(3.35), Inches(10.4), Inches(1.2),
         [(title, 44, True, WHITE)], line=1.08)
    rule(s, Inches(1.15), Inches(4.55), Inches(1.5))
    text(s, Inches(1.15), Inches(4.9), Inches(9.5), Inches(1.0),
         [(subtitle, 18, False, SAGE)], line=1.35)
    if note:
        notes(s, note)
    return s


def slide_content(title, blocks, kicker=None, note="", footer=None):
    s = prs.slides.add_slide(BLANK)
    bg(s, CREAM)
    band(s, 0, 0, W, Inches(0.14), FOREST)
    if kicker:
        text(s, Inches(0.9), Inches(0.62), Inches(11.5), Inches(0.35),
             [(kicker.upper(), 12, True, GOLD)])
        ty = Inches(1.0)
    else:
        ty = Inches(0.75)
    text(s, Inches(0.9), ty, Inches(11.5), Inches(1.0),
         [(title, 33, True, FOREST)], line=1.1)
    rule(s, Inches(0.9), ty + Inches(0.72), Inches(1.2))
    text(s, Inches(0.9), ty + Inches(1.05), Inches(11.5), Inches(4.9), blocks, line=1.3)
    if footer:
        text(s, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.4),
             [(footer, 11, False, SLATE)])
    if note:
        notes(s, note)
    return s


def slide_feature(num, name, tag, what, steps, say, extra=None, note_extra=None):
    s = prs.slides.add_slide(BLANK)
    bg(s, CREAM)
    band(s, 0, 0, W, Inches(0.14), FOREST)
    hdr = f"FEATURE {num}"
    if tag:
        hdr += f"   ·   {tag}"
    text(s, Inches(0.9), Inches(0.55), Inches(8.0), Inches(0.35),
         [(hdr, 12, True, GOLD)])
    text(s, Inches(0.9), Inches(0.92), Inches(8.4), Inches(0.85),
         [(name, 31, True, FOREST)], line=1.08)
    rule(s, Inches(0.9), Inches(1.78), Inches(1.2))
    text(s, Inches(0.9), Inches(2.05), Inches(7.3), Inches(0.9),
         [(what, 14, False, SLATE)], line=1.32)
    text(s, Inches(0.9), Inches(3.02), Inches(7.3), Inches(0.3),
         [("DEMO STEPS", 11, True, FOREST_MID)])
    runs = []
    for i, st in enumerate(steps):
        runs.append((f"{i+1}.  {st}", 14, False, INK, 7 if i else 0))
    text(s, Inches(0.9), Inches(3.38), Inches(7.3), Inches(3.3), runs, line=1.22)

    px, pw = Inches(8.65), Inches(3.78)
    band(s, px, Inches(2.05), pw, Inches(4.3), FOREST)
    band(s, px, Inches(2.05), Inches(0.06), Inches(4.3), GOLD)
    text(s, px + Inches(0.34), Inches(2.35), pw - Inches(0.68), Inches(0.3),
         [("SAY THIS", 11, True, GOLD_LIGHT)])
    text(s, px + Inches(0.34), Inches(2.75), pw - Inches(0.68), Inches(3.3),
         [(f"“{say}”", 14, False, WHITE)], line=1.34)
    if extra:
        text(s, Inches(0.9), Inches(6.72), Inches(7.3), Inches(0.5),
             [(extra, 11.5, True, GOLD)], line=1.2)
    body = f"SAY: {say}"
    if note_extra:
        body += f"\n\n{note_extra}"
    notes(s, body)
    return s


def slide_table(title, headers, rows, kicker=None, note="", col_w=None, fs=12):
    s = prs.slides.add_slide(BLANK)
    bg(s, CREAM)
    band(s, 0, 0, W, Inches(0.14), FOREST)
    if kicker:
        text(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.35),
             [(kicker.upper(), 12, True, GOLD)])
    text(s, Inches(0.9), Inches(0.92), Inches(11.5), Inches(0.8),
         [(title, 31, True, FOREST)], line=1.1)
    rule(s, Inches(0.9), Inches(1.72), Inches(1.2))
    n = len(rows) + 1
    tw = Inches(11.5)
    th = Inches(min(4.9, 0.5 * n))
    gt = s.shapes.add_table(n, len(headers), Inches(0.9), Inches(2.05), tw, th).table
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            gt.columns[i].width = Emu(int(tw * cw / total))
    for c, htxt in enumerate(headers):
        cell = gt.cell(0, c)
        cell.text = htxt
        cell.fill.solid(); cell.fill.fore_color.rgb = FOREST
        p = cell.text_frame.paragraphs[0]
        if p.runs:
            p.runs[0].font.size = Pt(12.5); p.runs[0].font.bold = True
            p.runs[0].font.color.rgb = WHITE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.12)
    for r, row in enumerate(rows, start=1):
        hero = str(row[0]).startswith("TerraVest")
        for c, val in enumerate(row):
            cell = gt.cell(r, c)
            cell.text = str(val)
            cell.fill.solid()
            if hero:
                cell.fill.fore_color.rgb = MIST
            else:
                cell.fill.fore_color.rgb = WHITE if r % 2 else SAGE
            p = cell.text_frame.paragraphs[0]
            if p.runs:
                p.runs[0].font.size = Pt(fs)
                p.runs[0].font.color.rgb = FOREST if hero else INK
                p.runs[0].font.bold = (c == 0) or hero
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.12)
    if note:
        notes(s, note)
    return s


def slide_matrix(title, cols, rows, kicker=None, note="", label_w=3.5, footer=None):
    """Capability matrix. Last column = TerraVest (highlighted). sym in {Y,P,N}"""
    SYM = {"Y": "●", "P": "◐", "N": "—"}
    s = prs.slides.add_slide(BLANK)
    bg(s, CREAM)
    band(s, 0, 0, W, Inches(0.14), FOREST)
    if kicker:
        text(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.35),
             [(kicker.upper(), 12, True, GOLD)])
    text(s, Inches(0.9), Inches(0.86), Inches(11.5), Inches(0.7),
         [(title, 29, True, FOREST)], line=1.1)
    rule(s, Inches(0.9), Inches(1.6), Inches(1.2))
    n = len(rows) + 1
    tw = Inches(11.5)
    th = Inches(min(4.4, 0.45 * n))
    tbl = s.shapes.add_table(n, len(cols) + 1, Inches(0.9), Inches(1.92), tw, th).table
    sym_w = (tw - Inches(label_w)) / len(cols)
    tbl.columns[0].width = Emu(int(Inches(label_w)))
    for i in range(1, len(cols) + 1):
        tbl.columns[i].width = Emu(int(sym_w))
    last = len(cols)
    for c, htxt in enumerate([""] + list(cols)):
        cell = tbl.cell(0, c)
        cell.text = htxt
        cell.fill.solid()
        cell.fill.fore_color.rgb = GOLD if c == last else FOREST
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
        if p.runs:
            p.runs[0].font.size = Pt(11)
            p.runs[0].font.bold = True
            p.runs[0].font.color.rgb = FOREST_DEEP if c == last else WHITE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.05)
    for r, (cap, syms) in enumerate(rows, start=1):
        cell = tbl.cell(r, 0)
        cell.text = cap
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if r % 2 else SAGE
        p = cell.text_frame.paragraphs[0]
        if p.runs:
            p.runs[0].font.size = Pt(11.5)
            p.runs[0].font.bold = True
            p.runs[0].font.color.rgb = INK
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.12)
        for c, sym in enumerate(syms, start=1):
            cell = tbl.cell(r, c)
            cell.text = SYM[sym]
            cell.fill.solid()
            cell.fill.fore_color.rgb = MIST if c == last else (WHITE if r % 2 else SAGE)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            if p.runs:
                p.runs[0].font.size = Pt(15 if sym != "N" else 13)
                p.runs[0].font.bold = True
                p.runs[0].font.color.rgb = (FOREST if sym == "Y"
                                            else GOLD if sym == "P" else GREY)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    text(s, Inches(0.9), Inches(1.92) + th + Inches(0.14), Inches(11.5), Inches(0.35),
         [("●  full capability          ◐  partial / workaround          "
           "—  not offered", 11, False, SLATE)])
    if footer:
        text(s, Inches(0.9), Inches(1.92) + th + Inches(0.46), Inches(11.5), Inches(0.7),
             [(footer, 12, True, FOREST)], line=1.25)
    if note:
        notes(s, note)
    return s


def slide_quote(quote, attrib=None, note=""):
    s = prs.slides.add_slide(BLANK)
    bg(s, FOREST_DEEP)
    band(s, 0, 0, Inches(0.28), H, GOLD)
    text(s, Inches(1.4), Inches(2.1), Inches(10.4), Inches(3.2),
         [(f"“{quote}”", 33, True, WHITE)], line=1.3, anchor=MSO_ANCHOR.MIDDLE)
    if attrib:
        rule(s, Inches(1.4), Inches(5.45), Inches(1.2))
        text(s, Inches(1.4), Inches(5.75), Inches(10.0), Inches(0.5),
             [(attrib, 14, False, GOLD_LIGHT)])
    if note:
        notes(s, note)
    return s


def slide_stats(title, stats, kicker=None, note=""):
    s = prs.slides.add_slide(BLANK)
    bg(s, CREAM)
    band(s, 0, 0, W, Inches(0.14), FOREST)
    if kicker:
        text(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.35),
             [(kicker.upper(), 12, True, GOLD)])
    text(s, Inches(0.9), Inches(0.92), Inches(11.5), Inches(0.8),
         [(title, 33, True, FOREST)], line=1.1)
    rule(s, Inches(0.9), Inches(1.72), Inches(1.2))
    cols, gap = 3, Inches(0.34)
    cw = (Inches(11.5) - gap * (cols - 1)) / cols
    ch = Inches(1.62)
    for i, (big, label) in enumerate(stats):
        r, c = divmod(i, cols)
        x = Inches(0.9) + c * (cw + gap)
        y = Inches(2.25) + r * (ch + Inches(0.3))
        band(s, x, y, cw, ch, WHITE)
        band(s, x, y, cw, Inches(0.055), GOLD)
        text(s, x + Inches(0.3), y + Inches(0.3), cw - Inches(0.6), Inches(0.7),
             [(big, 30, True, FOREST)], line=1.05)
        text(s, x + Inches(0.3), y + Inches(1.0), cw - Inches(0.6), Inches(0.5),
             [(label, 12, False, SLATE)], line=1.25)
    if note:
        notes(s, note)
    return s


# ================================================================ DECK
# ---------------------------------------------------------------- cover
slide_title(
    "All your wealth,\none place.",
    "Personal net worth, business finances, and property portfolio — in one login, "
    "for people who work for themselves.",
    kicker="TerraVest · Showcase & Demo Guide",
    footer="app.terravest.app   ·   Web · iOS · Android",
    note="DO NOT SHARE YOUR SCREEN YET.\n\n"
         "Deliver the problem first, then open the app. The dashboard is far more impressive to "
         "someone who has just been made to feel the mess it replaces.\n\n"
         "Full walkthrough ~25 min. Short version: features 1, 11, 12, 13, 14, 25, 26, 29-31, 34.")

slide_content(
    "What we'll cover",
    [("1.   The problem — and who has it", 17, True, INK),
     ("2.   The market — every competitor, what they charge, what they can't do", 17, True, INK, 12),
     ("3.   Personal money: one number, finally", 17, True, INK, 12),
     ("4.   The business command center", 17, True, INK, 12),
     ("5.   Property portfolio", 17, True, INK, 12),
     ("6.   Tax — where this gets different", 17, True, INK, 12),
     ("7.   The accountant handoff", 17, True, INK, 12),
     ("8.   Platform, trust, and the business model", 17, True, INK, 12)],
    kicker="Agenda",
    note="Keep this up for 15 seconds. It sets expectations and buys permission to move fast "
         "through the early sections.")

# ================================================================ PART 1
slide_section("01", "The Problem", "Why every existing money tool fails one specific person.")

slide_quote(
    "How much are you worth?",
    "The question Maria cannot answer — not because she's bad with money, but because the answer "
    "lives in four systems that don't talk.",
    note="THE OPENING STORY. Deliver slowly, no screen share yet.\n\n"
         "'Let me describe a real customer. Her name is Maria. She's a contractor.\n\n"
         "She has an LLC. She takes 1099 work. She owns three rental properties. And she files "
         "taxes jointly with her husband, who has a normal job.\n\n"
         "Now ask Maria a simple question. How much are you worth?\n\n"
         "She can't answer it. Not because she's bad with money - she's very good with money. But "
         "the answer is split across four systems that don't talk to each other. To find it, she "
         "needs an afternoon and a spreadsheet.\n\n"
         "And every March she spends a weekend collecting receipts, bank statements, and property "
         "numbers into a folder for her accountant - then pays that accountant professional rates "
         "to type it in.\n\n"
         "That's the problem. Not that the tools are bad. It's that no tool covers her whole "
         "life. So she does the joining by hand. Forever.'\n\nTHEN PAUSE.")

slide_content(
    "Three costs of the gap",
    [("No single number", 20, True, FOREST),
     ("She can't answer “what am I worth?” without an afternoon of work — so she never "
      "really knows, and can't make decisions from it.", 15, False, INK, 4),
     ("Tax money left on the table", 20, True, FOREST, 20),
     ("Rental depreciation, suspended passive losses, the 20% QBI deduction, quarterly estimates. "
      "Exactly the levers this customer has — and exactly the ones a consumer app never surfaces.",
      15, False, INK, 4),
     ("The accountant tax", 20, True, FOREST, 20),
     ("She pays professional rates for data assembly that software should do.", 15, False, INK, 4)],
    kicker="Why it matters",
    note="Use whichever of the three lands with this room. For users, costs 2 and 3 hit hardest. "
         "For investors, cost 1 frames the market.")

# ---------------------------------------------------------------- competition
slide_section("02", "The Market",
              "Every competitor, what they charge, and the rows none of them can fill.",
              note="This whole section is a credibility play. It proves you know your market "
                   "cold. Investors especially will test this.")

slide_table(
    "Every tool is built for someone who isn't her",
    ["Tool", "Price (list, Jul 2026)", "What it gives her", "What's still missing"],
    [["Monarch Money", "$14.99/mo · $99.99/yr\nPlus $199/yr",
      "Budgeting, net worth; Plus adds business & rental income tracking",
      "Income tracking only — no P&L, invoicing, AR, depreciation or QBI"],
     ["Copilot Money", "$13/mo · $95/yr",
      "Polished personal budgeting and investment tracking",
      "No business entity, no property, no tax modelling. No Android."],
     ["Empower", "Dashboard free\n0.89% AUM advisory",
      "Net worth and portfolio analysis, free",
      "Built to sell wealth management. No business books or rental accounting."],
     ["QuickBooks Online", "$20–$115/mo\n(Solopreneur → Plus)",
      "Full business accounting, invoicing, P&L",
      "No personal net worth, no investments, no property equity"],
     ["Stessa", "Free · Manage $12/mo\nPro $28/mo (annual)",
      "Rental income, expenses, Schedule E reporting",
      "Property only — no business entity, no personal wealth"],
     ["Baselane", "Free · Smart $20/mo",
      "Landlord banking, rent collection, bookkeeping",
      "Property only. Charges per lease, e-sign and screening."]],
    kicker="The gap — verified competitor pricing",
    col_w=[1.9, 2.1, 3.4, 4.1], fs=11,
    note="DON'T READ THE TABLE. Point at three rows:\n\n"
         "'Monarch is built for someone with a paycheck. QuickBooks is built for a bookkeeper. "
         "Stessa and Baselane only do property. Nobody covers all three.'\n\n"
         "IF ASKED ABOUT MONARCH PLUS: Monarch launched a $199/yr Plus tier adding 'business and "
         "rental income tracking'. Do NOT claim Monarch has no business features - that's now "
         "false and a knowledgeable investor will catch it. The honest, stronger answer:\n"
         "'Monarch Plus tracks business and rental INCOME as a category. It doesn't give you a "
         "P&L, invoicing, AR aging, multi-entity books, or depreciation and QBI in a tax engine. "
         "It categorises the money. We run the business.'\n\n"
         "Prices are list, verified July 2026. QuickBooks raised prices 1 Aug 2026 - re-check.")

slide_table(
    "The full market, and what each one charges",
    ["Product", "Price (list, Jul 2026)", "Who it's really for"],
    [["Monarch Money", "$14.99/mo · $99.99/yr · Plus $199/yr", "Households with a paycheck"],
     ["Copilot Money", "$13/mo · $95/yr", "Design-led personal budgeting (no Android)"],
     ["Empower", "Free dashboard · 0.89% AUM", "Lead-gen for wealth management"],
     ["YNAB", "$14.99/mo · $109/yr", "Zero-based budgeting discipline"],
     ["Quicken Simplifi", "$47.88/yr promo · $71.88 standard", "Low-cost personal budgeting"],
     ["QuickBooks Online", "$20 · $38 · $75 · $115/mo", "A bookkeeper doing the books"],
     ["Quicken Business & Personal", "$71.88/yr intro → $119.88 renewal",
      "Closest overlap — personal + business + rental"],
     ["Stessa", "Free · Manage $12 · Pro $28/mo", "Landlords, property only"],
     ["Baselane", "Free · Smart $20/mo", "Landlord banking + rent collection"],
     ["Hurdlr / Keeper", "$10–$20/mo · Keeper $199/yr", "Gig-worker mileage, deductions, filing"],
     ["TerraVest", "$9.99 Individual · $29.99 Business",
      "Self-employed with a business AND property"]],
    kicker="Competitive landscape",
    col_w=[3.0, 4.0, 4.5], fs=11,
    note="DON'T READ ALOUD. Point for five seconds and say:\n\n"
         "'I've mapped the whole market. Everyone here is excellent at one of the three things my "
         "customer needs. The one product that claims all three is Quicken - and that's the "
         "comparison I want to make.'\n\nThen move to the capability matrix.")

slide_matrix(
    "Personal and business capabilities",
    ["Monarch", "Copilot", "QuickBooks", "Stessa", "Quicken\nB&P", "TerraVest"],
    [("Personal net worth, all accounts",       ["Y", "Y", "N", "N", "Y", "Y"]),
     ("Budgets + debt payoff strategy lab",     ["P", "P", "N", "N", "Y", "Y"]),
     ("Multi-entity business books, isolated",  ["N", "N", "P", "N", "P", "Y"]),
     ("P&L, balance sheet, cash flow",          ["N", "N", "Y", "P", "Y", "Y"]),
     ("Invoicing + customer payment page",      ["N", "N", "Y", "N", "Y", "Y"]),
     ("AR aging + one-click bulk collections",  ["N", "N", "P", "N", "P", "Y"]),
     ("90-day cash forecast + shortfall alert", ["N", "N", "P", "N", "P", "Y"]),
     ("Vendor spend + contract renewal alerts", ["N", "N", "P", "N", "P", "Y"])],
    kicker="Capability matrix  ·  1 of 2",
    note="Read DOWN the TerraVest column, not across the rows.\n\n"
         "SAY: 'Monarch and Copilot own the top two rows. QuickBooks and Quicken own the middle. "
         "Nobody owns the column.'\n\n"
         "BE HONEST ABOUT QUICKEN. It genuinely does personal + business + rental in one product. "
         "Where we differ: isolated multi-entity books, AR aging as an analytic, and a real "
         "forward cash forecast - not a register with business categories.")

slide_matrix(
    "Property, tax and beyond",
    ["Monarch", "Copilot", "QuickBooks", "Stessa", "Quicken\nB&P", "TerraVest"],
    [("Property value, equity, cap rate",        ["P", "P", "N", "Y", "Y", "Y"]),
     ("Per-property expense tracking",           ["P", "N", "P", "Y", "Y", "Y"]),
     ("Rental depreciation + suspended losses",  ["N", "N", "P", "P", "P", "Y"]),
     ("QBI + proactive tax guidance engine",     ["N", "N", "N", "N", "N", "Y"]),
     ("Secure CPA share: expiry, passcode, log", ["N", "N", "P", "P", "N", "Y"]),
     ("AI assistant grounded in your own data",  ["Y", "P", "P", "N", "N", "Y"]),
     ("Real-estate deal + co-investment market", ["N", "N", "N", "N", "N", "Y"]),
     ("Web + iOS + Android",                     ["Y", "P", "Y", "Y", "Y", "Y"])],
    kicker="Capability matrix  ·  2 of 2",
    footer="Four rows nobody else can fill: depreciation with suspended losses · a QBI guidance "
           "engine · an auditable CPA share link · a deal marketplace.",
    note="THE SLIDE THAT WINS THE ARGUMENT. Four rows where the whole market is empty and we're "
         "full.\n\n"
         "SAY: 'Every one of these companies is good. But look at these four rows. Depreciation "
         "with carried-forward losses. A tax engine that tells her about QBI before April. A "
         "share link to her accountant that expires and logs who opened it. And a deal "
         "marketplace. Nobody else fills a single one of them.'\n\n"
         "Stessa/Quicken get partial credit on depreciation because they produce Schedule E "
         "reporting - they RECORD it, they don't model or advise on it. Say that if pressed; "
         "it's the accurate distinction and it protects you from looking like you're overclaiming.")

slide_content(
    "⚠  Quicken Business & Personal — the one that genuinely overlaps",
    [("It is the only product that does personal + business + rental in one subscription, with "
      "Schedule C/E reporting, invoicing and rental tracking — at about $120/yr on renewal, "
      "cheaper than TerraVest Business. Do not pretend it doesn't exist.", 15, False, INK),
     ("The honest answer — and the stronger one:", 15, True, FOREST, 14),
     ("“Quicken is the closest overlap and I take it seriously. But it's a register with "
      "business categories and Schedule C/E reporting — it records what already happened. We "
      "forecast what's coming, isolate books per entity, chase receivables, and advise on QBI and "
      "depreciation before April instead of reporting them after.”", 16, True, INK, 6),
     ("Where the difference is real and defensible:", 15, True, FOREST, 16),
     ("•  Isolated multi-entity books — Quicken is one file with categories, not separate businesses\n"
      "•  AR aging as an analytic, plus one-click bulk collections\n"
      "•  A genuine forward cash forecast with shortfall detection, not a register balance\n"
      "•  Proactive tax guidance (QBI, suspended losses) rather than after-the-fact Schedule E\n"
      "•  Auditable CPA sharing with expiry, passcode and access log\n"
      "•  An AI assistant grounded in the user's own numbers", 14, False, INK, 6)],
    kicker="Know this before you pitch",
    note="If an investor raises Quicken and you have no answer, it reads as not knowing your "
         "market. Now you have one.\n\n"
         "Volunteering this BEFORE they ask is even stronger - it signals you've done the work.")

slide_content(
    "What she pays today — and what she gets",
    [("To cover all three today, she buys three disconnected tools:", 15, False, SLATE),
     ("Monarch Core $14.99  +  QuickBooks Simple Start $38  +  Stessa Pro $35",
      19, True, INK, 12),
     ("≈ $88 / month  ·  three logins  ·  three data silos  ·  zero connection between them",
      17, True, GOLD, 8),
     ("TerraVest Business — $29.99/month, 7-day free trial", 22, True, FOREST, 22),
     ("One login. Personal net worth, multi-entity business books, and property — connected, with "
      "a tax engine that reads all three.", 16, False, INK, 6),
     ("About a third of the cost, and the only one where the three actually talk to each other.",
      16, True, FOREST, 14)],
    kicker="The stack she replaces",
    note="YOUR STRONGEST COMPETITIVE SLIDE. Concrete, verifiable, and it reframes price from a "
         "cost into a saving.\n\n"
         "SAY: 'Today she pays about eighty-eight dollars a month across three tools that don't "
         "talk to each other. We're thirty - and we're the only one where the three connect.'\n\n"
         "If challenged on the maths: list monthly prices as of July 2026. On annual billing the "
         "stack is about $74/month - still more than double us.\n\n"
         "Don't oversell the saving. The CONNECTION is the argument; price is the bonus.")

slide_content(
    "What TerraVest provides — precisely",
    [("Personal", 17, True, GOLD),
     ("Consolidated net worth from linked accounts · contributor breakdown · budgets with 50/30/20 "
      "· Avalanche/Snowball/Hybrid debt lab with extra-payment modelling · goals with "
      "required-monthly math · bill pay with duplicate protection", 13.5, False, INK, 3),
     ("Business — the part nobody else builds", 17, True, GOLD, 13),
     ("Isolated books per entity plus a consolidated view · business health score · 90-day cash "
      "forecast with shortfall detection · AR aging · invoicing with a public payment page and "
      "automatic transaction reconciliation · one-click chase of every overdue invoice · P&L, "
      "balance sheet and cash flow · budgets vs variance · cash-reserve and tax set-aside goals · "
      "vendor spend with renewal alerts · recurring-charge detection · QuickBooks import",
      13.5, False, INK, 3),
     ("Property and tax — the wedge", 17, True, GOLD, 13),
     ("Portfolio with value, equity and cap rate · per-property dated expense tracking · "
      "portfolio-wide tax export · rental depreciation via cost basis and land value · prior-year "
      "suspended loss carryforward · 20% QBI guidance · quarterly estimates · secure CPA sharing "
      "that is view-only, expiring, passcode-protected and access-logged",
      13.5, False, INK, 3),
     ("Platform", 17, True, GOLD, 13),
     ("AI assistant grounded in the user's real numbers · tamper-evident audit chain · encrypted "
      "identity data · export and delete · 9 languages · web, iOS and Android from one codebase",
      13.5, False, INK, 3)],
    kicker="Our offering",
    note="Your 'what do you actually do' slide. Don't read it - let them scan while you say:\n\n"
         "'Everything on this slide is built and running today. The business and property sections "
         "are the ones that don't exist anywhere else in one product.'\n\n"
         "If asked which is most used: the cash forecast and the tax engine.")

slide_quote(
    "TerraVest puts personal, business, and property in one place — and connects them.",
    "Because for a self-employed person, those were never three separate financial lives.",
    note="The pivot from problem to product. Now open the app.")

slide_stats(
    "What we've built",
    [("Live", "In production, HTTPS, deployed and monitored"),
     ("14", "Backend services"),
     ("305", "API endpoints"),
     ("35", "Product screens"),
     ("3", "Platforms — web, iOS, Android"),
     ("9", "Languages, incl. right-to-left Arabic"),
     ("174", "Database migrations"),
     ("82", "Data entities"),
     ("$9.99 / $29.99", "Individual / Business per month, 7-day trial")],
    kicker="Scale",
    note="For investors: this is not a prototype. It's live, deployed, monitored, and the revenue "
         "model is already built and running with feature gating live.\n\n"
         "Don't read every number - point at two or three.")

# ================================================================ PART 2
slide_section("03", "The Walkthrough",
              "46 features, in demo order. Personal → Business → Property → Tax → Trust → Platform.",
              note="From here, one slide per feature. Steps on the left, your line on the right "
                   "and in these notes.\n\nShort on time? Run 1, 11, 12, 13, 14, 25, 26, 29, 30, "
                   "31, 34.")

slide_content(
    "If you only have 10 minutes",
    [("Run these seven moments:", 16, False, SLATE),
     ("1   ·  Net Worth Dashboard — the one number", 17, True, INK, 14),
     ("11 ·  Multi-Business Command Center", 17, True, INK, 8),
     ("12 ·  90-Day Cash Forecast   ← strongest moment in the demo", 17, True, GOLD, 8),
     ("13 ·  AR Aging — the reason why", 17, True, INK, 8),
     ("14 ·  Invoicing & bulk reminders — the fix", 17, True, INK, 8),
     ("25 ·  Property portfolio", 17, True, INK, 8),
     ("29–31 ·  Tax engine   ← most differentiated screen", 17, True, GOLD, 8),
     ("34 ·  CPA secure sharing — the payoff", 17, True, INK, 8),
     ("If you can keep only two: the cash forecast and the tax screen. Nobody else in the market "
      "can show either.", 14, False, SLATE, 18)],
    kicker="Short version",
    note="Memorise this list. If the call runs short or the room gets restless, jump straight to "
         "feature 12.")

# ---------------- MODULE A
slide_section("A", "Personal Money",
              "Sidebar: Finance   ·   Features 1–10   ·   ~5 minutes",
              note="Play this module BRISK. It establishes competence but isn't your "
                   "differentiation. Only feature 1 gets real time.")

slide_feature(
    1, "Net Worth Dashboard", "★ CORE",
    "One number — everything owned minus everything owed — charted over time, with a breakdown of "
    "what moved it.",
    ["Land on Home (already logged in before the call).",
     "Pause on the net worth figure. Don't rush past it.",
     "Change the chart range to show it's live.",
     "Switch chart type — Area / Line / Bars.",
     "Open the “what moved it” contributor breakdown.",
     "Point out KPI cards and upcoming bills below."],
    "One number. Everything she owns minus everything she owes. This used to take her an "
    "afternoon. Now it's the home screen. And it doesn't just show the number — it shows why it "
    "changed.",
    note_extra="BONUS: if net worth drops >15% the chart shifts to a warning palette "
               "automatically. Mention it; don't try to trigger it live.")

slide_feature(
    2, "Account Linking", None,
    "Connect real bank and brokerage accounts. Balances and transactions flow in, and everything "
    "else updates from them.",
    ["Sidebar → Accounts.",
     "Show grouped account cards and the KPI row.",
     "Click Link account to open the connection flow.",
     "Select an institution; walk the bank's own login screen.",
     "Return to Accounts — show the new account in the list."],
    "This is how her money gets in. She picks her bank, logs in on the bank's own secure screen, "
    "and we never see her password. Balances and transactions flow from there.")

slide_feature(
    3, "Transactions", None,
    "Every transaction, filterable, with re-categorization that feeds budgets and taxes downstream.",
    ["Sidebar → Transactions.",
     "Filter by date range, then category, then amount.",
     "Sort a column.",
     "Re-categorize one transaction and show it update."],
    "Every transaction, filterable. And she can re-categorize anything — which matters, because "
    "these categories feed her budgets and her tax numbers later.")

slide_feature(
    4, "Budgets", None,
    "Spending plan with 50/30/20 presets across Needs, Wants and Savings, with over-budget alerts.",
    ["Sidebar → Budgets.",
     "Apply the 50/30/20 preset.",
     "Show Needs / Wants / Savings against actual spend.",
     "Change the period.",
     "Point out an over-budget alert."],
    "What she planned to spend, against what she actually spent — and it tells her when she's over.")

slide_feature(
    5, "Debt Lab", None,
    "Payoff strategy modelling — Avalanche vs Snowball vs Hybrid — with an extra-payment model and "
    "a recommendation.",
    ["Sidebar → Debt Lab.",
     "Show the debts table.",
     "Switch Avalanche → Snowball → Hybrid, reading each explainer.",
     "Enter an extra monthly payment; show the timeline change.",
     "Open the cross-strategy comparison — cheapest vs fastest.",
     "Point at the recommended plan."],
    "Avalanche is cheapest. Snowball is fastest to feel like progress. It compares them side by "
    "side, models paying extra, and recommends one. Most apps show you your debt. This one tells "
    "you what to do about it.")

slide_feature(
    6, "Goals", None,
    "Savings targets with progress tracking and required-monthly math.",
    ["Sidebar → Goals.",
     "Create a goal — name, target amount, target date.",
     "Show the required monthly contribution it calculates.",
     "Show progress on an existing goal."],
    "Set a target, and it works out exactly what she needs to put away each month to hit it.")

slide_feature(
    7, "Calculators", None,
    "Built-in financial math — mortgage payoff, extra payments, compound and simple interest.",
    ["Sidebar → Calculators.",
     "Run mortgage payoff with a realistic loan.",
     "Add an extra monthly payment — show interest saved and time cut.",
     "Show the compound interest calculator."],
    "The quick math, built in — so she isn't opening a browser tab mid-decision.")

slide_feature(
    8, "Investments", None,
    "Holdings, connected brokers, and alternative investments.",
    ["Sidebar → Investments.",
     "Show the holdings list.",
     "Show connect a broker.",
     "Open Alternatives and add an alternative investment."],
    "Her investments — including alternatives, which matters for this customer, because a lot of "
    "their wealth isn't in a normal brokerage account.",
    extra="[MISSING: confirm Alternatives add-flow fields; is Marketplace demo-ready?]")

slide_feature(
    9, "Cash Position", None,
    "What's actually liquid right now, with account freshness status.",
    ["Sidebar → Cash.",
     "Show the cash position total.",
     "Point at status indicators — Healthy / Stale / Action required."],
    "What she actually has liquid — and it flags accounts that have gone stale, so the number "
    "stays trustworthy.",
    extra="[MISSING: confirm layout — do card balances live here or under My Business?]")

slide_feature(
    10, "Bill Pay", None,
    "Multi-step payment: payee → amount → funding → schedule → confirm, with duplicate protection "
    "and cancellation.",
    ["Sidebar → Pay Bills.",
     "Select or add a payee.",
     "Enter an amount; choose the funding account.",
     "Pick a date, then confirm — show it appear in the list.",
     "Cancel the scheduled payment to show it's reversible.",
     "Point at the sidebar badge showing pending count."],
    "Who she's paying, how much, from which account, and when. It's built to be safe — "
    "double-click and it won't pay twice, and she can cancel before it goes out.",
    note_extra="Say 'scheduled' or 'submitted' - not 'the money has gone'.")

# ---------------- MODULE B
slide_section("B", "The Business Command Center",
              "Sidebar: My Business   ·   Features 11–24   ·   ~7 minutes",
              note="SLOW DOWN. Your strongest module; features 12-14 are the strongest sequence in "
                   "the entire demo.\n\n"
                   "TRANSITION IN: 'That's her personal money - and honestly, other apps do a "
                   "version of that. Here's what nobody else does.'\n\n"
                   "WARNING: My Business is behind a FeatureGate. Confirm your demo account "
                   "renders the real page, not an upgrade prompt.")

slide_feature(
    11, "Multi-Business Command Center", "★★ CORE",
    "Each company as its own isolated set of books, plus a combined view, with a health score "
    "across the top.",
    ["Sidebar → My Business → Overview tab.",
     "Show the business switcher — move between entities.",
     "Show the consolidated view across all businesses.",
     "Point at the business health score."],
    "Here's where we're different from every other money app. This is her LLC as its own set of "
    "books, completely separate from her personal money. Three companies means three sets of "
    "books, plus a combined view.",
    note_extra="TRANSITION IN: 'That's her personal money - and other apps do a version of that. "
               "Here's what nobody else does.'")

slide_feature(
    12, "90-Day Cash Forecast", "★★ CORE — STRONGEST MOMENT",
    "Forward-looking cash projection that detects a shortfall before it happens.",
    ["On Overview, scroll to the 90-day cash forecast.",
     "Trace the projection line to the shortfall point.",
     "STOP TALKING for two full seconds. Let them read it."],
    "This is a 90-day cash forecast. It's telling her that in about nine weeks, she runs tight. "
    "Most business owners find out they have a cash problem when the money is already gone. This "
    "tells her while she can still do something about it. That's the difference between a "
    "dashboard and a decision.",
    note_extra="THE SINGLE MOST IMPORTANT MOMENT IN THE DEMO.\n\n"
               "After 'she runs tight in nine weeks', STOP TALKING FOR TWO SECONDS. The silence "
               "does the work. Business owners have a physical reaction to this screen.\n\n"
               "Rehearse this out loud twice before the call.")

slide_feature(
    13, "AR Aging", "★ CORE",
    "Everyone who owes money, bucketed by how overdue they are.",
    ["From the forecast, click into AR aging.",
     "Show the aging buckets — 30 / 60 / 90+ days.",
     "Point at the most overdue invoice."],
    "And it shows her why. These are her unpaid invoices, sorted by how late they are. This one is "
    "60 days overdue. That's her cash problem, sitting right there.",
    note_extra="TRANSITION: 'And it doesn't just warn her. It shows her WHY.'\n\n"
               "Needs genuinely overdue invoices in the demo account or there's nothing here.")

slide_feature(
    14, "Invoicing & Bulk Reminders", "★ CORE",
    "Create and send invoices, give customers a public payment page, reconcile payments to bank "
    "transactions, and chase every overdue invoice at once.",
    ["My Business → Business Tools → Invoices.",
     "Create an invoice — customer, line items, amount, due date.",
     "Send it by email; show the confirmation.",
     "Open the public invoice page in a second tab, as the customer sees it.",
     "Return and click “Remind all overdue”.",
     "Show payment reconciliation linked to a transaction."],
    "One click just chased every overdue invoice. Her customer gets a link, opens a payment page, "
    "and pays. When the money lands, it matches back against her bank transaction automatically. "
    "Problem, reason, and fix — one screen, three clicks. Not a report. A decision.",
    note_extra="Have a second browser profile ready to show the customer's view of the invoice "
               "page. Far more convincing than describing it.")

slide_feature(
    15, "Business Transactions Ledger", None,
    "Period-aware ledger merging linked bank data with manual entries.",
    ["My Business → Transactions tab.",
     "Change the period; show figures update.",
     "Add a manual transaction.",
     "Show linked and manual entries in one ledger."],
    "Her business ledger — bank data and anything she enters by hand, in one place, for whatever "
    "period she's looking at.")

slide_feature(
    16, "Business Expenses", None,
    "Dated, categorized business expenses linked to real transactions, with export.",
    ["My Business → Expenses tab.",
     "Add an expense — date, category, amount, vendor.",
     "Link it to a transaction.",
     "Export to CSV."],
    "Every business expense, dated, categorized, and tied back to the real transaction — so at tax "
    "time there's nothing to reconstruct.")

slide_feature(
    17, "Business Credit Cards", None,
    "Business card tracking, kept separate from personal cards.",
    ["My Business → Credit Cards tab.",
     "Show card balances and activity."],
    "Business cards tracked separately from personal — which sounds obvious, and is exactly what "
    "gets messy when you're running both from one wallet.",
    extra="[MISSING: are cards added manually or pulled from linked accounts?]")

slide_feature(
    18, "Budgets & Variance", None,
    "Business spending plan versus actuals.",
    ["My Business → Business Tools → Budgets.",
     "Set a budget for a category.",
     "Show the variance against actual spend."],
    "What she planned to spend in the business, against what she actually spent.")

slide_feature(
    19, "Business Goals", None,
    "Cash reserve and tax set-aside targets.",
    ["My Business → Business Tools → Goals.",
     "Set a cash reserve target.",
     "Set a tax set-aside target.",
     "Show progress against both."],
    "A cash reserve, and a tax set-aside — so the tax money is put away before she spends it. "
    "Every self-employed person has learned that lesson the hard way.",
    note_extra="The 'learned that lesson the hard way' line gets a laugh of recognition from "
               "self-employed audiences. Pause after it.")

slide_feature(
    20, "Vendor Management", None,
    "Vendor spend, contract status, and renewal alerts.",
    ["My Business → Business Tools → Vendors.",
     "Show computed vendor spend, ranked.",
     "Open a vendor; set status, renewal date, notes.",
     "Point at a renewal alert."],
    "Everyone she pays, what she pays them, and a warning before a contract auto-renews.")

slide_feature(
    21, "Recurring & Subscription Detection", None,
    "Automatic detection of repeating charges, plus per-customer payment behaviour.",
    ["My Business → Business Tools → Recurring.",
     "Show detected recurring charges.",
     "Show per-customer payment behaviour — who pays on time."],
    "It spots subscriptions and repeat charges by itself. And on the income side, it learns which "
    "customers actually pay on time — which is worth knowing before you take the next job.")

slide_feature(
    22, "Business Reports", None,
    "Cash-basis P&L, balance sheet, and cash flow — exportable.",
    ["My Business → Reports tab.",
     "Open the Profit & Loss.",
     "Open the Balance Sheet.",
     "Open the Cash Flow statement.",
     "Export to CSV, then show print-to-PDF."],
    "Profit and loss, balance sheet, cash flow — the three statements her accountant will ask for, "
    "exportable as a spreadsheet or a PDF.")

slide_feature(
    23, "Business Documents", None,
    "Per-business document folders and invoice attachments.",
    ["My Business → Documents tab.",
     "Upload a file to a business folder.",
     "Show an invoice with an attachment.",
     "Show the All view across businesses, newest first."],
    "Per-business folders, and receipts attached to the invoices they belong to.")

slide_feature(
    24, "QuickBooks Import", None,
    "Import existing books from QuickBooks Online.",
    ["My Business → Connect QuickBooks.",
     "Show the connection prompt."],
    "And if she already keeps books in QuickBooks, it imports — she isn't starting over.",
    extra="[MISSING: demo steps for the post-connection import view]")

# ---------------- MODULE C
slide_section("C", "Property Portfolio",
              "Sidebar: Real Estate   ·   Features 25–28   ·   ~3 minutes",
              note="TRANSITION IN: 'That's the business. Now the third piece - the properties.'\n\n"
                   "Steady pace. This module sets up the tax module that follows.")

slide_feature(
    25, "Property Portfolio", "★ CORE",
    "Each property with estimated value, mortgage balance, real equity, and rental cap rate.",
    ["Sidebar → Properties.",
     "Show the portfolio — value, debt, equity per property.",
     "Add a property — address, purchase price, mortgage.",
     "Show the auto-estimated value from market data.",
     "Point at the cap rate for a rental."],
    "Her three rentals. What each is worth, what she still owes, and the equity she actually has — "
    "plus the yield, so she can see which property is genuinely earning.")

slide_feature(
    26, "Per-Property Expense Tracker", "★ CORE",
    "Dated, categorized expenses per property — the shoebox, digitized.",
    ["From Properties, open a property → Expenses drawer.",
     "Add an expense — date, category, amount.",
     "Show the category breakdown for the year.",
     "Export the property's expenses.",
     "Show the portfolio-wide combined tax export."],
    "Every repair, insurance bill, property tax payment, and management fee — dated and "
    "categorized, all year long. This is the part that usually lives in a shoebox.",
    note_extra="TRANSITION OUT: 'Now - why does all this property detail matter so much? Because "
               "of what it unlocks.' Then go to Tax.")

slide_feature(
    27, "Deal Room", None,
    "A marketplace of sponsor real-estate deals, with track records.",
    ["Sidebar → Deal Room.",
     "Browse the deal list.",
     "Open a deal — terms, documents, sponsor detail.",
     "Show the sponsor track record.",
     "Watch a deal, then express interest."],
    "When she wants to put money to work — deals from sponsors, with the sponsor's track record "
    "visible before she trusts them.",
    note_extra="FOR INVESTORS ADD: 'This is also a future revenue line - we sit between sponsors "
               "and qualified investors whose financial picture we already know.'")

slide_feature(
    28, "Fractional LLC", None,
    "Co-investment marketplace — buy into a property alongside other investors.",
    ["Sidebar → Fractional LLC.",
     "Show available co-investment opportunities.",
     "Open one; show the ownership structure."],
    "And she can co-invest in a property alongside other investors, instead of having to buy one "
    "outright.",
    extra="[MISSING: demo steps for the participation / commit flow]")

# ---------------- MODULE D
slide_section("D", "Tax & Documents",
              "Sidebar: Taxes, Documents   ·   Features 29–35   ·   ~6 minutes",
              note="YOUR DIFFERENTIATION. Second-slowest module after Business.\n\n"
                   "TRANSITION IN: 'Now - why does all that property detail matter so much? "
                   "Because of what it unlocks.'\n\n"
                   "Feature 30 is the one nobody else can show. Do not rush it.")

slide_feature(
    29, "Tax Estimator", "★★ CORE",
    "A live, year-round tax picture: filing status, W-2 income, self-employment income, investment "
    "income, and rentals.",
    ["Sidebar → Taxes.",
     "Set filing status — choose Married filing jointly.",
     "Add the spouse's W-2.",
     "Enter self-employment / 1099 income.",
     "Show interest, dividends, capital gains, retirement income.",
     "Show the running estimate updating as you type."],
    "This is her tax picture — live all year, not just in April. Her filing status, her husband's "
    "W-2, her 1099 income from the business.")

slide_feature(
    30, "Rental Tax Modelling", "★★ CORE — MOST DIFFERENTIATED",
    "The rental section no consumer app has: depreciation and suspended loss carryforward.",
    ["On Taxes, open the rental section.",
     "Enter gross rents collected.",
     "Enter mortgage interest, property tax, insurance, repairs, management.",
     "Enter property cost basis and land value.",
     "Enter prior-year suspended loss (carryforward).",
     "Show the resulting deduction change in the estimate."],
    "Cost basis and land value — that's depreciation. Every rental owner can deduct part of the "
    "building's value every single year. And prior-year suspended loss — when a rental loses "
    "money, that loss doesn't disappear. It carries forward. Most people forget it exists and "
    "never claim it.",
    note_extra="SLOW RIGHT DOWN HERE. Explain what each term IS - most people in the room won't "
               "know, and teaching them something is what makes the product feel expert rather "
               "than pretty.\n\nOne of the three moments to rehearse out loud.")

slide_feature(
    31, "Tax Guidance Engine", "★ CORE",
    "Rule-based prompts for the deductions and credits this customer qualifies for.",
    ["On Taxes, scroll to the guidance panel.",
     "Show Deductions — rental depreciation, suspended losses, 20% QBI, tax-advantaged accounts.",
     "Show Credits — child & dependent care, home energy and EV.",
     "Show the quarterly estimated payments prompt."],
    "It tells her what she's missing. The 20% QBI deduction on her business income. Her quarterly "
    "estimates. Energy credits. A normal budgeting app has no idea any of these exist — and these "
    "are exactly the deductions this customer has.",
    note_extra="THE LINE THAT LANDS THE MODULE:\n\n'This screen is the difference between an app "
               "that watches your money and an app that UNDERSTANDS your money.'")

slide_feature(
    32, "Document Upload & OCR", None,
    "Upload tax documents; the app extracts values from PDFs and photos.",
    ["On Taxes, use the document upload.",
     "Upload a W-2 (PDF or photo).",
     "Show extracted values populating the form.",
     "Show remove / replace on an uploaded file."],
    "She can photograph a W-2 instead of typing it in.",
    extra="[MISSING: pre-test OCR on your demo file — do not attempt this live untested]")

slide_feature(
    33, "Document Center", "★ CORE",
    "Personal document storage with folders and a cross-app registry pulling in documents "
    "generated elsewhere in the product.",
    ["Sidebar → Documents.",
     "Show the folder structure.",
     "Upload a document to a folder.",
     "Show the cross-app registry — Real Estate and Business exports appearing here."],
    "Every financial document in one place — including the exports the app generates for her "
    "automatically.",
    note_extra="TRANSITION IN: 'And here's how the year ends.'")

slide_feature(
    34, "CPA Secure Sharing", "★ CORE — THE PAYOFF",
    "Share with an accountant via a view-only link with expiry, mandatory passcode, and a full "
    "access log.",
    ["From Documents, select files → Share.",
     "Build a multi-file share set.",
     "Set an expiry date.",
     "Set the passcode (mandatory).",
     "Generate the link; open it in an incognito window as the accountant.",
     "Enter the passcode; show the view-only experience.",
     "Return and show the access log — who opened it, when."],
    "Remember Maria's weekend of collecting paperwork? This is that — one link. View-only. It "
    "expires. It needs a passcode. And she can see exactly when her accountant opened it. Her "
    "whole tax year, sent safely, in about ten seconds.",
    note_extra="THE EMOTIONAL PAYOFF OF THE WHOLE DEMO. Tie it explicitly back to the "
               "weekend-with-a-shoebox image from your opening.\n\n"
               "Third of the three moments to rehearse. Have the incognito window pre-opened.")

slide_feature(
    35, "CPA Marketplace", None,
    "Find and connect with an accountant.",
    ["Navigate to /cpa (route-only — not in the sidebar).",
     "Browse available CPAs.",
     "Open a profile; show the connect flow."],
    "And if she doesn't have an accountant, she can find one here.",
    extra="[MISSING: confirm the in-app entry point so you're not typing a URL live]")

# ---------------- MODULE E
slide_section("E", "Intelligence & Trust",
              "Features 36–40   ·   ~2 minutes",
              note="Quick pace. Shows completeness and reassures on security. Feature 39 is strong "
                   "with technical investors.")

slide_feature(
    36, "AI Assistant", None,
    "A financial assistant grounded in the user's actual accounts, business and properties — with "
    "scope control, response styles, a prompt library, and voice input.",
    ["Sidebar → AI Assistant.",
     "Show the insight cards on arrival.",
     "Ask ONE pre-tested question.",
     "Show the scope selector — what it's allowed to see.",
     "Show response styles and the prompt library.",
     "Point at the disclaimer."],
    "An assistant that can actually see her numbers — not generic advice from the internet. She "
    "controls what it's allowed to look at.",
    extra="⚠  Test your exact question before the call. Never ask an AI live what you haven't "
          "seen answer well.")

slide_feature(
    37, "Notifications & Messages", None,
    "In-app inbox plus email / SMS / push alerts with per-user preference control.",
    ["Click the notifications bell in the topbar.",
     "Click View all messages → Messages page.",
     "Settings → Notifications; toggle a preference.",
     "Send a test notification."],
    "Alerts in one inbox, and she controls exactly which ones reach her and how.")

slide_feature(
    38, "Security Center", None,
    "Two-factor authentication, active session management, and a readable login history with "
    "device and location.",
    ["Sidebar → Security.",
     "Show two-factor authentication settings.",
     "Show active sessions.",
     "Show login history — friendly labels, device, time, location."],
    "Two-factor, every active session, and a full login history — device, time, and where it came "
    "from. If something looks wrong, she sees it.")

slide_feature(
    39, "Audit Trail", None,
    "A tamper-evident hash chain of every state change, with an endpoint that proves the log "
    "wasn't altered.",
    ["Show the user-facing activity timeline.",
     "For technical audiences: open /api/v1/audit/verify and show the chain verifying."],
    "Every change to an account is written into a cryptographic chain. This endpoint re-walks that "
    "chain and proves nothing was altered or deleted. Most fintechs retrofit this after their "
    "first audit. We built it before we had users.",
    note_extra="STRONG WITH TECHNICAL INVESTORS. A credibility moment - it signals you think about "
               "trust architecturally, not as a feature.")

slide_feature(
    40, "Data Export & Account Deletion", None,
    "Full data export and self-service account deletion.",
    ["Settings → Data & Privacy.",
     "Trigger a data export.",
     "Show the delete account option — describe it, don't click it."],
    "She can export everything or delete her account outright. Her data is hers. We say that, and "
    "then we make it a button.")

# ---------------- MODULE F
slide_section("F", "Platform & Business Model",
              "Features 41–46   ·   ~2 minutes",
              note="Fast. For investors, close on feature 43 (Subscription) - your revenue proof "
                   "point.")

slide_feature(
    41, "Profile & Identity", None,
    "Personal details and identity information, with SSN/EIN encrypted and shown only as last four "
    "digits.",
    ["Sidebar footer → Profile.",
     "Show name and contact details.",
     "Show the masked SSN / EIN — last four only.",
     "Show notification preference toggles."],
    "Her identity details are encrypted and only ever displayed as the last four digits — even to "
    "her.")

slide_feature(
    42, "Settings, Themes & Languages", None,
    "Appearance, regional settings, and full internationalization.",
    ["Sidebar → Settings.",
     "Cycle the theme from the topbar — Light / Dark / Glass.",
     "Change the language; show a right-to-left language such as Arabic.",
     "Show regional settings."],
    "Light, dark, and glass. Nine languages, including right-to-left for Arabic — because this "
    "customer base isn't only English-speaking.",
    note_extra="The theme switch is a nice visual flourish. Two seconds, always gets a reaction.")

slide_feature(
    43, "Subscription & Plans", "INVESTOR CLOSE",
    "Two-tier subscription with a 7-day trial and live feature gating — built, running, and "
    "repriceable without a deploy.",
    ["Sidebar → Subscription.",
     "Show Individual $9.99/mo and Business $29.99/mo.",
     "Show the 7-day trial and annual pricing (~2 months free).",
     "Show the per-plan feature list.",
     "Mention the trial-ending banner and onboarding modal."],
    "The plans, the trial, the billing lifecycle, and the feature locks are all built and running. "
    "And pricing lives in the database, not the code — we can reprice or repackage without a "
    "deploy.",
    note_extra="YOUR INVESTOR CLOSE. Follow immediately with the 'why can't the incumbents do "
               "this' beat.\n\nAlso the natural place to reprise the stack-cost argument: she "
               "replaces ~$88/mo of disconnected tools with $29.99.")

slide_feature(
    44, "Learn & Guide", None,
    "In-app educational modules.",
    ["Topbar help icon → Learn.",
     "Show the educational modules."],
    "Built-in education, because half of this customer's problem is that nobody ever explained "
    "depreciation to them.",
    extra="[MISSING: which Learn modules have finished content worth showing?]")

slide_feature(
    45, "Mobile Apps", None,
    "The same product on iPhone and Android from one codebase, installable as a PWA.",
    ["Show the app on a phone if one is to hand.",
     "Otherwise show the install-to-home-screen prompt in the browser."],
    "Same product on iPhone and Android, from one codebase. She can install it from the browser.")

slide_feature(
    46, "Ops Portal", "INTERNAL — INVESTORS ONLY",
    "A staff support console with identities fully separate from customer identities, "
    "permission-based access control, caller verification, and an actor/target audit chain.",
    ["Describe rather than demo — separate login at /ops.",
     "If showing: customer 360, caller verification, tiered disclosure.",
     "Show PII reveal requiring a typed reason."],
    "Support staff have completely separate accounts — a staff token is refused on customer "
    "routes, so an agent's credential can never act as a customer. Revealing a customer's personal "
    "data requires a typed reason that gets permanently recorded.",
    extra="[MISSING: decide whether to show at all — needs a seeded staff account]")

# ================================================================ PART 3
slide_section("04", "Presenter Notes",
              "Transitions, pacing, the three moments, and the questions you'll be asked.")

slide_table(
    "Transitions carry the narrative",
    ["Between", "Say"],
    [["Open → Personal", "“Let me show you what Maria sees on a Monday morning.”"],
     ["Personal → Business",
      "“Other apps do a version of that. Here's what nobody else does.”"],
     ["Forecast → AR aging", "“And it doesn't just warn her. It shows her why.”"],
     ["AR aging → Invoicing", "“And she fixes it from the same screen.”"],
     ["Business → Property",
      "“That's the business. Now the third piece — the properties.”"],
     ["Property → Tax",
      "“Why does this matter so much? Because of what it unlocks.”"],
     ["Tax → Documents", "“And here's how the year ends.”"],
     ["Documents → Close", "“One link. Instead of a folder and a weekend.”"]],
    kicker="Script notes",
    col_w=[3.0, 8.5],
    note="Say these deliberately. They're what makes it a story rather than a tour.")

slide_table(
    "Pacing",
    ["Module", "Time", "How to play it"],
    [["A · Personal", "5 min", "Brisk. Establishes competence. Only feature 1 gets real time."],
     ["B · Business", "7 min", "Slowest section. Features 12–14 are your strongest sequence."],
     ["C · Property", "3 min", "Steady. Sets up tax."],
     ["D · Tax & Docs", "6 min", "Second-slowest. Feature 30 is your differentiation."],
     ["E · Intelligence", "2 min", "Quick. Shows completeness."],
     ["F · Platform", "2 min", "Fast. Close on 43 for investors."]],
    kicker="Script notes",
    col_w=[2.6, 1.4, 7.5],
    note="Total ~25 minutes. If running long, compress A, E and F - never B or D.")

slide_content(
    "The three moments that decide the demo",
    [("12  ·  The cash forecast", 21, True, GOLD),
     ("After you say “she runs tight in nine weeks,” stop talking for two full seconds. "
      "The silence does the work. Business owners have a physical reaction to this screen.",
      15, False, INK, 4),
     ("30  ·  Depreciation and suspended losses", 21, True, GOLD, 18),
     ("Slow right down and explain what each one is. Most people in the room won't know — and "
      "teaching them something is what makes the product feel expert rather than pretty.",
      15, False, INK, 4),
     ("34  ·  The CPA link", 21, True, GOLD, 18),
     ("The emotional payoff. Tie it explicitly back to the weekend-with-a-shoebox image from your "
      "opening.", 15, False, INK, 4),
     ("Rehearse these three out loud, twice. Everything else you can narrate naturally from the "
      "app in front of you.", 14, True, FOREST, 20)],
    kicker="Script notes",
    note="Don't memorise 46 features. Memorise these three.")

slide_content(
    "Closing for USERS",
    [("They care about one thing: does this solve my problem?", 16, False, SLATE),
     ("“You don't need to change how you work. Link your accounts, add your properties once, "
      "and it keeps itself current. The tax deductions alone — the depreciation, the "
      "carried-forward losses, the 20% business deduction — usually cover the cost of the app many "
      "times over.”", 17, True, INK, 14),
     ("Then turn it into a conversation:", 15, True, FOREST, 18),
     ("•  How do you track your rentals today?", 15, False, INK, 8),
     ("•  How long does tax preparation take you each year?", 15, False, INK, 4),
     ("•  What's the most annoying part of your money every month?", 15, False, INK, 4),
     ("•  If you could delete one tool you use today, which one?", 15, False, INK, 4),
     ("Write their answers down during the call. They're worth more than the rest of the demo.",
      14, True, GOLD, 16)],
    kicker="Audience: end users",
    note="A user demo should END WITH YOU LISTENING. The last five minutes are research, not "
         "presentation.")

slide_content(
    "Closing for INVESTORS",
    [("The question they're already thinking:", 15, True, FOREST),
     ("“The obvious question is why Mint or QuickBooks doesn't just build this. They can. But "
      "it costs them their focus. If Monarch adds a P&L, it makes their product worse for the 95% "
      "of users who have a normal paycheck. If QuickBooks adds personal net worth, it confuses the "
      "bookkeeper they sell to. Serving this customer means combining three worlds none of them "
      "wants to own all three of. That gap is our opening.”", 16, True, INK, 10),
     ("Then the trajectory:", 15, True, FOREST, 16),
     ("“Today it's a trusted dashboard. Next is advice. After that, a financial product — "
      "lending or advisory — sold to people whose complete financial picture only we can see. What "
      "I'm chasing isn't more features. It's 50 to 100 weekly active users in this niche.”",
      16, True, INK, 10)],
    kicker="Audience: investors",
    note="Deliver immediately after showing the Subscription screen (feature 43).")

slide_content(
    "The close",
    [("Return to Home. Ending where you started closes the loop visually.", 15, False, SLATE),
     ("“So — one login. Her personal net worth. Her business, with a warning before cash runs "
      "out. Her properties. Her taxes, with the deductions she actually qualifies for. And one "
      "safe link to her accountant.", 19, True, INK, 14),
     ("For millions of people whose money life is too complicated for a budgeting app — and too "
      "personal for accounting software.”", 19, True, INK, 10),
     ("Then stop talking. Silence invites questions. Filling it undoes the ending.",
      15, True, GOLD, 18)],
    kicker="Script notes",
    note="Do not add features after this. The ending only works if you let it end.")

slide_table(
    "Questions you'll be asked",
    ["Question", "Answer"],
    [["“Isn't this just Monarch?”",
      "“Monarch's excellent for someone with a paycheck. Their new $199 Plus tier does add "
      "business and rental income tracking — but that's categorising income. No P&L, no invoicing, "
      "no AR aging, no multi-entity books, no depreciation or QBI. It watches the money. We run "
      "the business.”"],
     ["“What about QuickBooks?”",
      "“It answers ‘how is the business doing?’ It has no idea what its user is "
      "worth — no personal net worth, investments or property equity. And it's $38 to $115 a month "
      "for the tiers this customer needs.”"],
     ["“What about Quicken Business & Personal?”",
      "“The closest overlap, and I take it seriously — personal, business and rental in one, "
      "about $120/yr on renewal. It's a register with business categories and Schedule C/E "
      "reporting. It records what happened. We forecast what's coming, isolate books per entity, "
      "chase receivables, and advise on QBI and depreciation before April.”"],
     ["“How big is the market?”",
      "“Tens of millions of self-employed Americans, a large share owning rental property. "
      "But I'm not making the case on market size — I'm making it on traction in the niche.”"],
     ["“Do you have users?”",
      "“Not yet at scale — the product went live recently. That's the next milestone. I'd "
      "rather give you the honest number than a vanity one.”"]],
    kicker="Q&A",
    col_w=[2.8, 8.7], fs=11,
    note="Answer 'do you have users' with the honest number. Hedging costs more credibility than "
         "the number does.")

slide_table(
    "Harder questions",
    ["Question", "Answer"],
    [["“What's the moat?”",
      "“Three layers. Near term, focus — incumbents can't follow without damaging their core "
      "product. Medium term, data — once someone's personal, business and property history lives "
      "here, switching means rebuilding years of records. Long term, the financial product.”"],
     ["“How do you handle security?”",
      "“JWT enforced independently in every service. SSN/EIN encrypted, last-four only. "
      "Secrets in a KMS-backed store. Services refuse to boot on a weak secret. Every state change "
      "in a tamper-evident hash chain. Staff access fully separated from customer identity.”"],
     ["“Are you SOC 2 compliant?”",
      "“No formal certification — that's an audit with real cost, premature before revenue. "
      "But I built the primitives an audit asks for, because they're expensive to retrofit. It's a "
      "process to run, not an architecture to rebuild.”"],
     ["“Why so many services?”",
      "“Partly over-engineered, and I'll own that. Each domain has a different external "
      "dependency and compliance surface. The cost is operational overhead, which I control by "
      "running Compose on one VM rather than Kubernetes.”"],
     ["“What's the biggest risk?”",
      "“Distribution, not product. I've built something this customer needs. I haven't yet "
      "proven I can reach them efficiently. That's why the milestone is weekly actives.”"]],
    kicker="Q&A",
    col_w=[2.8, 8.7], fs=11,
    note="Owning the over-engineering point and the SOC 2 gap BEFORE they press reads as judgment. "
         "Being caught not having considered them does not.")

slide_content(
    "Lines worth memorizing",
    [("“All your wealth, one place — personal, business, and property.”",
      19, True, FOREST),
     ("“Every money app assumes your finances are simple. If you work for yourself, they "
      "never are.”", 19, True, FOREST, 14),
     ("“Mint doesn't know you own a business. QuickBooks doesn't know you own a house. Nobody "
      "knows you own both.”", 19, True, FOREST, 14),
     ("“It warns you about a cash problem while you can still fix it. That's the difference "
      "between a dashboard and a decision.”", 19, True, FOREST, 14),
     ("“This is the difference between an app that watches your money and an app that "
      "understands it.”", 19, True, FOREST, 14),
     ("“She replaces about eighty-eight dollars a month of disconnected tools with thirty.”",
      19, True, FOREST, 14)],
    kicker="Script notes",
    note="Six lines. Learn these and you can improvise everything else.")

# ================================================================ PART 4
slide_section("05", "Running the Demo Well",
              "Six things that will help more than rehearsing feature 37.")

slide_content(
    "Setup that protects you",
    [("Use two screens", 20, True, FOREST),
     ("Share only the app window; keep this deck in presenter view on the other. Never scroll a "
      "script in a shared window — it instantly reveals you're reading.", 15, False, INK, 4),
     ("Record a backup video", 20, True, FOREST, 16),
     ("Record the full walkthrough once, cleanly, and keep it in a background tab. If the site is "
      "slow or Wi-Fi drops: “let me show you the recorded version while that loads.” "
      "Demos are lost in dead air, not to bugs.", 15, False, INK, 4),
     ("Rehearse three moments, not forty-six", 20, True, FOREST, 16),
     ("Features 12, 30 and 34 — twice, out loud. Everything else you can narrate naturally from "
      "the app in front of you.", 15, False, INK, 4)],
    kicker="Preparation",
    note="The two-screen setup is the single highest-value change you can make. No second monitor? "
         "Print the feature slides. Paper doesn't get shared by accident.")

slide_content(
    "⚠  Seed the demo account — and check the feature gate",
    [("Several of your best features have nothing to show without data:", 15, False, SLATE),
     ("•  Feature 12 cash forecast — needs transaction history to project", 15, False, INK, 10),
     ("•  Feature 13 AR aging — needs genuinely overdue invoices", 15, False, INK, 4),
     ("•  Feature 14 bulk remind — needs multiple overdue invoices", 15, False, INK, 4),
     ("•  Feature 20 vendors — needs recurring vendor payments", 15, False, INK, 4),
     ("•  Feature 26 property expenses — needs a year of dated expenses", 15, False, INK, 4),
     ("•  Features 29–31 tax — needs income, W-2 and rental figures", 15, False, INK, 4),
     ("CRITICAL: My Business, Deal Room and Fractional LLC sit behind a feature gate. On the wrong "
      "plan or an expired trial, your flagship screen renders an UPGRADE PROMPT instead of the "
      "product. Log in the day before and confirm all three.", 16, True, GOLD, 18)],
    kicker="Preparation",
    note="The item most likely to ruin a call. Verify the gate the day before, not an hour before.")

slide_content(
    "Demo hygiene",
    [("•  Log in before the call and leave the tab open — never demo through a login or an OTP "
      "code.", 15, False, INK),
     ("•  Hard-refresh (Cmd+Shift+R). This is a PWA; a stale cached bundle can load.",
      15, False, INK, 8),
     ("•  Close every other tab. Turn off system and browser notifications.", 15, False, INK, 8),
     ("•  Zoom the browser to ~110% — financial tables are unreadable shrunk down.",
      15, False, INK, 8),
     ("•  Pre-open your closing tab (Subscription for investors).", 15, False, INK, 8),
     ("•  Click the whole path once the day before, so nothing surprises you live.",
      15, False, INK, 8),
     ("•  Have a second browser profile ready for the public invoice page and the CPA share link "
      "— showing the recipient's view is far more convincing than describing it.",
      15, False, INK, 8),
     ("Handling questions mid-demo: “Great question — that's actually the next screen.” "
      "Or park it, note it, and genuinely return before you close.", 14, True, FOREST, 16)],
    kicker="Preparation",
    note="Not returning to a parked question is worse than not taking it.")

slide_table(
    "Still to fill in",
    ["#", "Feature", "What's needed"],
    [["8", "Investments", "Confirm Alternatives add-flow; is Marketplace demo-ready?"],
     ["9", "Cash Position", "Confirm layout — do card balances live here?"],
     ["17", "Credit Cards", "Manual entry or pulled from linked accounts?"],
     ["24", "QuickBooks", "Demo steps for the post-connection import view"],
     ["28", "Fractional LLC", "Demo steps for the participation / commit flow"],
     ["32", "Document OCR", "Pre-test extraction on your actual demo file"],
     ["35", "CPA Marketplace", "Confirm the in-app entry point (not in sidebar)"],
     ["44", "Learn & Guide", "Which modules have finished content?"],
     ["46", "Ops Portal", "Show it at all? Needs a seeded staff account"]],
    kicker="Open items",
    col_w=[0.8, 2.6, 8.1], fs=11.5,
    note="Also verify before any demo:\n"
         "- My Business / Deal Room / Fractional LLC render past the FeatureGate\n"
         "- The AI Assistant answers your exact planned question well\n"
         "- The demo account has overdue invoices, or features 13 and 14 show nothing\n"
         "- Competitor prices are still current (QuickBooks raised prices 1 Aug 2026)")

# ---- cheat card
s = prs.slides.add_slide(BLANK)
bg(s, FOREST_DEEP)
band(s, 0, 0, Inches(0.28), H, GOLD)
text(s, Inches(1.1), Inches(0.72), Inches(11.0), Inches(0.5),
     [("ONE-PAGE CHEAT CARD", 13, True, GOLD_LIGHT)])
text(s, Inches(1.1), Inches(1.18), Inches(11.0), Inches(0.7),
     [("Keep this visible during the call", 30, True, WHITE)], line=1.1)
rule(s, Inches(1.1), Inches(1.95), Inches(1.2))
text(s, Inches(1.1), Inches(2.3), Inches(5.3), Inches(4.4),
     [("Live", 12, True, GOLD), ("app.terravest.app", 16, True, WHITE, 2),
      ("For", 12, True, GOLD, 12),
      ("Self-employed people with a business and property", 16, True, WHITE, 2),
      ("Pricing", 12, True, GOLD, 12),
      ("Individual $9.99 · Business $29.99 · 7-day trial", 16, True, WHITE, 2),
      ("She replaces", 12, True, GOLD, 12),
      ("~$88/mo of disconnected tools (Monarch + QuickBooks + Stessa)", 16, True, WHITE, 2),
      ("Scale", 12, True, GOLD, 12),
      ("14 services · 305 endpoints · 35 screens · 3 platforms", 16, True, WHITE, 2)],
     line=1.18)
text(s, Inches(7.0), Inches(2.3), Inches(5.3), Inches(4.4),
     [("Order", 12, True, GOLD),
      ("Home → Business (forecast → AR → remind) → Property → Tax "
       "→ Documents → Close", 16, True, WHITE, 2),
      ("Best moment", 12, True, GOLD, 12),
      ("90-day cash forecast", 16, True, WHITE, 2),
      ("Best screen", 12, True, GOLD, 12),
      ("Tax — depreciation, suspended losses, QBI", 16, True, WHITE, 2),
      ("Closest competitor", 12, True, GOLD, 12),
      ("Quicken B&P ~$120/yr — records; we forecast & advise", 16, True, WHITE, 2),
      ("Biggest risk (own it)", 12, True, GOLD, 12),
      ("Distribution, not product", 16, True, WHITE, 2)],
     line=1.18)
notes(s, "Print this slide. It's everything you might blank on.")

slide_quote(
    "Too complicated for a budgeting app. Too personal for accounting software.",
    "TerraVest  ·  app.terravest.app",
    note="Leave this on screen while you take questions.")

out = "/Users/suresh/Desktop/projects/my-wealth-management/DOCUMENTATION/TerraVest-Showcase-Deck.pptx"
prs.save(out)
print(f"Saved: {out}")
